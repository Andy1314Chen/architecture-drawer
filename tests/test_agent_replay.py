"""Opt-in real-agent replay: each eval's ``input.md`` -> a real coding agent
authors ``gen.py``; the harness is the deterministic score + artifact gate.

This is the replay layer closest to real usage: the skill is *installed*
into a leak-free sandbox, and a real agent harness (the **Pi coding agent**,
pi.dev) discovers it, reads ``input.md``, and self-authors the generator —
exactly what a user does in their editor. The
harness then re-runs the produced ``gen.py`` deterministically (no trust in
agent self-reports) and asserts:

  - the final quality score clears ``AGENT_REPLAY_MIN_SCORE`` (80);
  - the SVG exists and is XML-parseable;
  - the PPTX exists, opens via ``python-pptx``, and has shapes;
  - the PNG exists and is non-empty (tolerated per-case when ``rsvg-convert``
    is absent — PNG rasterization is a soft dep).

Anti-leakage: only the skill's public surface (no ``evals/``) + the target
eval's ``input.md`` enter the sandbox. The golden ``gen.py`` is never copied,
so the agent cannot transcribe it.

Skipped by default. Run with ``pytest --agent-replay [--agent-iter N]
[--agent-eval <name>]``. Needs the ``pi`` CLI and a configured provider
(``ANTHROPIC_API_KEY`` or a prior ``pi /login``). Local/nightly only.
"""
from __future__ import annotations

import shutil
import sys
import xml.etree.ElementTree as ET
from pathlib import Path

import pytest

from agent_backends import PiAgentBackend, prepare_sandbox
from conftest import ROOT, SCRIPTS, AGENT_REPLAY_MIN_SCORE, score_gen_script

if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))
from semantic_qa import run_semantic_qa          # noqa: E402

# --- agent prompts ---------------------------------------------------------
# Short, skill-discovery-respecting instructions. They deliberately do NOT
# inline SKILL.md: the whole point is that the agent loads the installed skill
# itself. They also never start with ``-`` (avoids argv mis-parse as a flag).

_AGENT_INITIAL = """\
Load the `architecture-drawer` skill, then read `input.md` in this directory.
Following the skill's Generate-Evaluate-Correct workflow, write a complete,
self-contained `gen.py` that draws the architecture described in `input.md`
using the SVGDrawer DSL.

Resolve the skill's scripts/ directory as documented in `AGENTS.md` (this
sandbox), import from `svg_utils`, `evaluator`, `svg2pptx`, `design_brief`,
and `semantic_qa`, and:

- declare a `BRIEF = DesignBrief(...)` (Step 1 design brief: layout, flow,
  palette_role keyed by data-node-id, flow_chain for pipeline stages) and
  give the tinted layer containers matching `node_id=` values;
- call `evaluate_svg(drawer)` and print `print(f"Score: {score}")`;
- call `run_semantic_qa(drawer, expected_size=(W, H), brief=BRIEF)` and print
  its report — no `brief-*` FAIL/WARN may remain;
- iterate (run the script, read its report, adjust coordinates) until the score
  is 100 with no `[FAIL]` items — use `auto_refine` for geometric fixes and fix
  dangles / crossings / text-overlaps by moving nodes and rerouting edges;
- save the artifact triplet plus `brief.json` (via `BRIEF.write(...)`) next to
  `gen.py`: SVG via `save_svg`, PNG via `rasterize_svg`, PPTX via
  `svg2pptx.svg_to_pptx`;
- write only `gen.py` and its artifacts; do not edit `input.md` or the skill.

When done, print the final score line.
"""

_AGENT_REFINE = """\
If `gen.py` is missing or does not run, create or fix it first — follow the
architecture-drawer skill and the spec in `input.md`. Otherwise `gen.py` still
has evaluator `[FAIL]` or `[WARN]` items, or scores below 100: re-read it, run
it (`python3 gen.py`) to see the current score and report, then fix the issues
by adjusting coordinates, layout, and edge routing — do NOT add suppressions or
silence checks. Re-run until the score is 100 with no `[FAIL]`. Re-export the
SVG/PNG/PPTX triplet next to `gen.py`, then print the final score line.
"""


# --- deterministic gate helpers -------------------------------------------

def _has_fail(report: str) -> bool:
    return any("[FAIL]" in ln for ln in report.splitlines())


def _score(sandbox: Path) -> tuple[int | None, str]:
    gen = sandbox / "gen.py"
    if not gen.is_file():
        return None, "no gen.py produced"
    return score_gen_script(sandbox)


def _latest(paths: list[Path]) -> Path | None:
    """Most-recently-modified file (by mtime); None when empty. The agent may
    write several artifacts during a loop — pick the newest, not the
    alphabetically-last filename."""
    return max(paths, key=lambda p: p.stat().st_mtime, default=None)


def _check_artifacts(sandbox: Path) -> list[str]:
    """Validate the produced SVG/PPTX/PNG/brief.json. PNG is optional when
    rsvg is absent; the brief contract (brief.json) is mandatory — an agent
    that skips the design brief fails the replay."""
    problems = []
    svg = _latest(list(sandbox.glob("*.svg")))
    if svg is None:
        return ["no SVG produced"]
    try:
        ET.parse(str(svg))
    except ET.ParseError as exc:
        problems.append(f"SVG does not parse: {exc}")
    pptx = _latest(list(sandbox.glob("*.pptx")))
    if pptx is None:
        problems.append("no PPTX produced")
    else:
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx))
            if not prs.slides or not prs.slides[0].shapes:
                problems.append("PPTX opens but has no shapes")
        except Exception as exc:  # noqa: BLE001 - surface any pptx failure
            problems.append(f"PPTX does not open: {exc}")
    # PNG needs rsvg-convert; tolerate its absence per-case.
    png = _latest(list(sandbox.glob("*.png")))
    if (png is None or png.stat().st_size == 0) and shutil.which("rsvg-convert"):
        problems.append("no/empty PNG produced (rsvg-convert is installed)")
    brief = _latest(list(sandbox.glob("brief.json")))
    if brief is None:
        problems.append("no brief.json produced (design-brief contract skipped)")
    else:
        try:
            sys.path.insert(0, str(SCRIPTS))
            from design_brief import DesignBrief  # noqa: WPS433 (sandboxed)
            DesignBrief.load(str(brief))
        except Exception as exc:  # noqa: BLE001 - any parse failure fails
            problems.append(f"brief.json does not parse: {exc}")
    return problems


def _persist_sandbox(sandbox: Path, eval_name: str, score: int | None, report: str) -> Path:
    """Copy a sandbox's run into ``output/agent_replay/<eval_name>/`` for review.

    Keeps ``gen.py``, the artifact triplet, a copy of ``input.md``, the
    installed skill (under ``.pi/skills/``), and writes the final score +
    evaluator report as ``score_report.txt``. Returns the destination dir.
    ``output/`` is gitignored, so this never enters version control.
    """
    dst = ROOT / "output" / "agent_replay" / eval_name
    if dst.exists():
        shutil.rmtree(dst, ignore_errors=True)
    shutil.copytree(sandbox, dst, dirs_exist_ok=True)
    (dst / "score_report.txt").write_text(
        f"eval: {eval_name}\nscore: {score}\n\n{report}\n", encoding="utf-8"
    )
    return dst


def _semantic_gate(sandbox: Path) -> tuple[list[str], list[str]]:
    """Run semantic QA on the produced SVG against the eval's input.md spec.

    Returns (fail_lines, warn_lines). FAILs gate the run (regenerate / fix);
    WARNs are fed back into the next refine round so the agent can address
    them (e.g. rails slicing containers, missing spec entities).
    """
    svg = _latest(list(sandbox.glob("*.svg")))
    if svg is None:
        return [], []
    spec = sandbox / "input.md"
    try:
        qa = run_semantic_qa(
            svg.read_text(encoding="utf-8"),
            spec_text=spec.read_text(encoding="utf-8") if spec.is_file() else None,
        )
    except Exception:  # noqa: BLE001 - semantic QA must never crash the gate
        return [], []
    fails = [i.render() for i in qa.issues if i.severity == "fail"]
    warns = [i.render() for i in qa.issues if i.severity != "fail"]
    return fails, warns


def _agent_replay_one(
    backend: PiAgentBackend, eval_dir: Path, max_iter: int, *,
    name: str, keep: bool = False,
) -> tuple[int, list[str]]:
    """Run one eval through the agent loop; return (best_score, problems).

    Round 0 sends the initial prompt. Then up to *max_iter* stateless refine
    rounds: each re-scores the on-disk ``gen.py`` and, if below 100 or with
    [FAIL] items, asks the agent to fix it. It also stops early once a
    non-FAIL score past ``AGENT_REPLAY_MIN_SCORE`` plateaus (no gain over the
    prior round), so a stable ~95 does not burn every remaining refine call.
    Finally the artifacts are gated.

    If *keep* is set, the sandbox (gen.py + SVG/PNG/PPTX + skill + final score
    report) is copied to ``output/agent_replay/<name>/`` instead of deleted.
    """
    sandbox = prepare_sandbox(eval_dir)
    best = 0
    problems: list[str] = []
    last_score: int | None = None
    last_report = ""
    prev_score: int | None = None
    try:
        backend.run(_AGENT_INITIAL, sandbox)                 # round 0
        score, report = _score(sandbox)
        last_score, last_report = score, report
        best = max(best, score or 0)
        for _ in range(max_iter):
            if score is not None and score >= 100 and not _has_fail(report):
                break
            # Early exit: already past the floor with no FAIL and the last
            # refine round made no progress — each extra round is a full Pi
            # call, so stop rather than chase a 100 it isn't reaching.
            if (
                prev_score is not None
                and score is not None
                and score >= AGENT_REPLAY_MIN_SCORE
                and not _has_fail(report)
                and score <= prev_score
            ):
                break
            prev_score = score
            # Semantic feedback: feed this round's semantic-QA findings into
            # the (stateless) refine prompt so the agent fixes meaning-level
            # defects — rails over components, missing spec entities —
            # alongside the geometric ones.
            _, sem_warns = _semantic_gate(sandbox)
            refine = _AGENT_REFINE
            if sem_warns:
                refine += ("\n\nSemantic QA findings from semantic_qa.py — fix "
                           "these too (reroute rails off components, add "
                           "missing spec component labels):\n"
                           + "\n".join(sem_warns[:10]))
            backend.run(refine, sandbox)              # stateless refine
            score, report = _score(sandbox)
            last_score, last_report = score, report
            best = max(best, score or 0)
        problems += _check_artifacts(sandbox)
        sem_fails, _ = _semantic_gate(sandbox)
        problems += sem_fails
        return best, problems
    finally:
        if keep:
            try:
                _persist_sandbox(sandbox, name, last_score, last_report)
            except Exception as exc:  # noqa: BLE001 - never block cleanup
                # Persisting is an opt-in debug aid; a failure here must not
                # mask the real result or leak the sandbox via the rmtree
                # that would otherwise be skipped.
                print(
                    f"[agent-replay] warning: could not persist {name}: {exc}",
                    file=sys.stderr,
                )
        shutil.rmtree(sandbox, ignore_errors=True)


# --- the test --------------------------------------------------------------

@pytest.mark.agent_replay
def test_agent_replay_quality(eval_cases, request):
    """Real-agent replay: install skill, let pi author gen.py, gate on score."""
    if not request.config.getoption("--agent-replay"):
        pytest.skip("--agent-replay not passed; deterministic gate only")
    if not PiAgentBackend.available():
        pytest.skip(
            "'pi' CLI not on PATH; install @earendil-works/pi-coding-agent to "
            "run agent replay, or omit --agent-replay for the deterministic gate."
        )

    max_iter = request.config.getoption("--agent-iter")
    only = request.config.getoption("--agent-eval")
    keep = request.config.getoption("--agent-keep")
    candidates = {
        n: d for n, d in eval_cases.items() if (d / "input.md").is_file()
    }
    if only:
        matched = {n: d for n, d in candidates.items() if n == only}
        if not matched:
            pytest.fail(
                f"--agent-eval {only!r} matched no eval case with an "
                f"input.md spec; available: {', '.join(sorted(candidates))}"
            )
        candidates = matched
    if not candidates:
        pytest.skip("no eval with an input.md spec to replay")

    backend = PiAgentBackend(
        provider=request.config.getoption("--agent-provider"),
        model=request.config.getoption("--agent-model"),
    )
    failures: list[str] = []
    checked = 0
    for name, eval_dir in sorted(candidates.items()):
        checked += 1
        try:
            best, problems = _agent_replay_one(
                backend, eval_dir, max_iter, name=name, keep=keep)
        except Exception as exc:  # noqa: BLE001 - isolate one case's crash
            failures.append(f"[{name}] agent run error: {exc}")
            continue
        if best < AGENT_REPLAY_MIN_SCORE:
            failures.append(
                f"[{name}] agent best score {best} < floor {AGENT_REPLAY_MIN_SCORE}"
            )
        for p in problems:
            failures.append(f"[{name}] {p}")

    if failures:
        pytest.fail(
            f"agent replay failed {len(failures)} check(s) across {checked} "
            f"case(s): " + "; ".join(failures)
        )
