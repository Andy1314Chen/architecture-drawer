# -*- coding: utf-8 -*-
"""
SVG to PowerPoint (PPTX) export — native editable shapes.

Converts SVG content into native, editable PowerPoint shapes: rectangles,
ovals, connectors, text boxes, and freeforms. Inspired by and modeled after
the svg2pptx project (github.com/benouinirachid/svg2pptx), but self-contained
and tuned for diagrams produced by SVGDrawer.

Two export modes:
  - "shapes" (default): each SVG element becomes an editable PPTX shape.
    Best for diagrams you want to tweak in PowerPoint/Keynote.
  - "image": rasterize the SVG to PNG via rsvg-convert and embed as a picture.
    Perfect visual fidelity (arrows, curves, everything), but NOT individually
    editable. Use when the SVG is complex and you just need it to look right.

Coordinate system: SVG pixels are mapped to PowerPoint EMU (1 px = 9525 EMU at
96 DPI). The SVG viewBox is scaled to fit the slide while preserving aspect
ratio (letterboxed, centered).

Usage:
    from svg2pptx import svg_to_pptx

    # From a file
    svg_to_pptx("diagram.svg", "diagram.pptx")

    # From an SVG string (what SVGDrawer.render() returns)
    svg_to_pptx(drawer.render(), "diagram.pptx")

    # Image mode (raster fallback for complex SVGs)
    svg_to_pptx("diagram.svg", "diagram.pptx", mode="image")

    # Add to an existing presentation's slide
    from svg2pptx import add_svg_to_slide
    add_svg_to_slide(svg_string, slide)
"""
from __future__ import annotations

import html
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional, Union
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# Reuse the battle-tested affine transform helpers from svg_utils (same dir).
from svg_utils import (IDENTITY, multiply_matrix, parse_transform,
                       transform_point)

# ---------------------------------------------------------------------------
# Units
# ---------------------------------------------------------------------------
# 1 inch = 914400 EMU, 1 inch = 96 px  ->  1 px = 9525 EMU.
EMU_PER_PX = 9525


def px(v: float) -> int:
    """Convert SVG pixels to EMU (rounded)."""
    return int(round(v * EMU_PER_PX))


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
@dataclass
class PptxConfig:
    """Export configuration.

    slide_w / slide_h : slide dimensions in inches (default 13.333 x 7.5 = 16:9).
    margin            : letterbox margin in inches.
    scale             : extra scale multiplier on top of fit-to-slide.
    curve_tolerance   : Bezier flattening tolerance in px (lower = smoother).
    mode              : "shapes" (native editable) or "image" (rasterized).
    image_dpi         : DPI for raster mode.
    default_fill      : fill used when SVG omits it ("none" or hex).
    default_stroke    : stroke used when SVG omits it.
    """

    slide_w: float = 13.333
    slide_h: float = 7.5
    margin: float = 0.0
    scale: float = 1.0
    curve_tolerance: float = 1.0
    mode: str = "shapes"
    image_dpi: int = 200
    default_fill: str = "none"
    default_stroke: str = "none"


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------
_NAMED = {
    "white": "#FFFFFF", "black": "#000000", "red": "#FF0000",
    "green": "#008000", "blue": "#0000FF", "yellow": "#FFFF00",
    "none": "none", "transparent": "none",
}


def _normalize_color(val: Optional[str], default: str = "none") -> str:
    """Normalize a CSS/SVG color to '#rrggbb' (lowercase) or 'none'."""
    if not val or val.strip() == "":
        return default
    val = val.strip().lower()
    val = _NAMED.get(val, val)
    if val.startswith("#"):
        h = val[1:]
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        if len(h) == 6:
            return f"#{h}"
    return val if val == "none" else default


def _hex_to_rgbcolor(hex_color: str) -> Optional[RGBColor]:
    """Convert '#rrggbb' to RGBColor, or None for 'none'."""
    hex_color = _normalize_color(hex_color)
    if hex_color == "none":
        return None
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Style parsing (presentation attributes + style="..." attribute)
# ---------------------------------------------------------------------------
@dataclass
class SvgStyle:
    fill: str = "none"
    stroke: str = "none"
    stroke_width: float = 1.0
    fill_opacity: float = 1.0
    stroke_opacity: float = 1.0
    opacity: float = 1.0
    stroke_dasharray: str = ""
    font_size: float = 14.0
    font_weight: str = "normal"
    font_style: str = "normal"
    font_family: str = "Arial, sans-serif"
    text_anchor: str = "start"
    dominant_baseline: str = "auto"


_STYLE_PROPS = {
    "fill", "stroke", "stroke-width", "fill-opacity", "stroke-opacity",
    "opacity", "stroke-dasharray", "font-size", "font-weight", "font-style",
    "font-family", "text-anchor", "dominant-baseline",
}


def _parse_style(el: ET.Element, parent: Optional[SvgStyle] = None) -> SvgStyle:
    """Merge presentation attributes, inline style, and parent inheritance."""
    base = parent if parent else SvgStyle()
    s = SvgStyle(
        fill=base.fill, stroke=base.stroke, stroke_width=base.stroke_width,
        fill_opacity=base.fill_opacity, stroke_opacity=base.stroke_opacity,
        opacity=base.opacity, stroke_dasharray=base.stroke_dasharray,
        font_size=base.font_size, font_weight=base.font_weight,
        font_style=base.font_style, font_family=base.font_family,
        text_anchor=base.text_anchor, dominant_baseline=base.dominant_baseline,
    )
    # Inline style attribute
    inline = {}
    style_attr = el.get("style", "")
    if style_attr:
        for decl in style_attr.split(";"):
            decl = decl.strip()
            if ":" in decl:
                k, v = decl.split(":", 1)
                inline[k.strip()] = v.strip()
    # Apply — inline style overrides presentation attrs in SVG spec, but for
    # diagrams generated programmatically the attrs are the source of truth.
    for prop in _STYLE_PROPS:
        val = el.get(prop)
        if val is None:
            val = inline.get(prop)
        if val is None:
            continue
        val = val.strip()
        if prop == "fill":
            s.fill = _normalize_color(val, "none")
        elif prop == "stroke":
            s.stroke = _normalize_color(val, "none")
        elif prop == "stroke-width":
            s.stroke_width = _float(val, s.stroke_width)
        elif prop == "fill-opacity":
            s.fill_opacity = _float(val, s.fill_opacity)
        elif prop == "stroke-opacity":
            s.stroke_opacity = _float(val, s.stroke_opacity)
        elif prop == "opacity":
            s.opacity = _float(val, s.opacity)
        elif prop == "stroke-dasharray":
            s.stroke_dasharray = val
        elif prop == "font-size":
            s.font_size = _float(val, s.font_size)
        elif prop == "font-weight":
            s.font_weight = val
        elif prop == "font-style":
            s.font_style = val
        elif prop == "font-family":
            s.font_family = val
        elif prop == "text-anchor":
            s.text_anchor = val
        elif prop == "dominant-baseline":
            s.dominant_baseline = val
    return s


def _effective_opacity(style: SvgStyle) -> float:
    return min(1.0, max(0.0,
                        style.opacity * style.fill_opacity))


def _effective_stroke_opacity(style: SvgStyle) -> float:
    return min(1.0, max(0.0, style.opacity * style.stroke_opacity))


def _float(val: str, default: float = 0.0) -> float:
    try:
        return float(re.sub(r"[^0-9.\-]", "", val))
    except (ValueError, TypeError):
        return default


# ---------------------------------------------------------------------------
# SVG path parsing + Bezier flattening
# ---------------------------------------------------------------------------
_PATH_CMD = re.compile(r"([MLHVCSTAQZmlhvcstaqz])")
_NUM = re.compile(
    r"[-+]?(?:\d+\.\d*|\.\d+|\d+)(?:[eE][-+]?\d+)?"
)


def _tokenize_path(d: str):
    """Yield (command, [floats...]) tuples from an SVG path 'd' attribute."""
    # Split into command + numbers segments.
    parts = _PATH_CMD.split(d)
    # parts[0] is leading whitespace/garbage before first command
    for i in range(1, len(parts), 2):
        cmd = parts[i]
        nums = [float(m.group()) for m in _NUM.finditer(parts[i + 1])]
        yield cmd, nums


def _flatten_path(d: str, tol: float = 1.0):
    """Flatten an SVG path into a list of sub-paths (polylines).

    Returns: list of (points, closed) where points is [(x,y), ...].
    Curves (C/S/Q/T/A) are subdivided into line segments within `tol` px.
    """
    subpaths = []
    points: list[tuple[float, float]] = []
    closed = False
    cx = cy = 0.0  # current point
    sx = sy = 0.0  # subpath start point

    def moveto(x, y, relative=False):
        nonlocal cx, cy, sx, sy
        if relative:
            x, y = cx + x, cy + y
        cx, cy = x, y
        sx, sy = x, y

    def lineto(x, y, relative=False):
        nonlocal cx, cy
        if relative:
            x, y = cx + x, cy + y
        points.append((x, y))
        cx, cy = x, y

    for cmd, nums in _tokenize_path(d):
        it = iter(nums)

        if cmd in ("M", "m"):
            if points:
                subpaths.append((points, closed))
            x, y = next(it), next(it)
            moveto(x, y, cmd == "m")
            points = [(cx, cy)]
            closed = False
            # subsequent pairs are implicit lineto
            for px2, py2 in _pairs(it):
                lineto(px2, py2, cmd == "m")

        elif cmd in ("L", "l"):
            for x, y in _pairs(it):
                lineto(x, y, cmd == "l")

        elif cmd in ("H", "h"):
            for x in it:
                lineto(x if cmd == "H" else cx + x, cy)

        elif cmd in ("V", "v"):
            for y in it:
                lineto(cx, y if cmd == "V" else cy + y)

        elif cmd in ("C", "c"):
            for x1, y1, x2, y2, x3, y3 in _group(it, 6):
                if cmd == "c":
                    x1, y1 = cx + x1, cy + y1
                    x2, y2 = cx + x2, cy + y2
                    x3, y3 = cx + x3, cy + y3
                _flatten_cubic(cx, cy, x1, y1, x2, y2, x3, y3, tol, points)
                cx, cy = x3, y3

        elif cmd in ("S", "s"):
            # smooth cubic — reflect previous control point
            for x2, y2, x3, y3 in _group(it, 4):
                if cmd == "s":
                    x2, y2 = cx + x2, cy + y2
                    x3, y3 = cx + x3, cy + y3
                x1 = 2 * cx - _last_ctrl[0]
                y1 = 2 * cy - _last_ctrl[1]
                _flatten_cubic(cx, cy, x1, y1, x2, y2, x3, y3, tol, points)
                _last_ctrl[0], _last_ctrl[1] = x2, y2
                cx, cy = x3, y3

        elif cmd in ("Q", "q"):
            for x1, y1, x2, y2 in _group(it, 4):
                if cmd == "q":
                    x1, y1 = cx + x1, cy + y1
                    x2, y2 = cx + x2, cy + y2
                _flatten_quad(cx, cy, x1, y1, x2, y2, tol, points)
                _last_ctrl[0], _last_ctrl[1] = x1, y1
                cx, cy = x2, y2

        elif cmd in ("T", "t"):
            for x2, y2 in _pairs(it):
                if cmd == "t":
                    x2, y2 = cx + x2, cy + y2
                x1 = 2 * cx - _last_ctrl[0]
                y1 = 2 * cy - _last_ctrl[1]
                _flatten_quad(cx, cy, x1, y1, x2, y2, tol, points)
                _last_ctrl[0], _last_ctrl[1] = x1, y1
                cx, cy = x2, y2

        elif cmd in ("A", "a"):
            for rx_, ry_, angle, large, sweep, x, y in _group(it, 7):
                if cmd == "a":
                    x, y = cx + x, cy + y
                _arc_to_cubics(cx, cy, rx_, ry_, angle, bool(large),
                               bool(sweep), x, y, tol, points)
                cx, cy = x, y

        elif cmd in ("Z", "z"):
            closed = True
            cx, cy = sx, sy

    if points:
        subpaths.append((points, closed))
    return subpaths


_last_ctrl = [0.0, 0.0]  # last Bezier control point (for S/T smooth commands)


def _pairs(it):
    """Yield (x, y) pairs from a flat iterator."""
    while True:
        try:
            x = next(it)
            y = next(it)
            yield x, y
        except StopIteration:
            return


def _group(it, n):
    """Yield n-tuples from a flat iterator."""
    while True:
        chunk = []
        try:
            for _ in range(n):
                chunk.append(next(it))
        except StopIteration:
            if len(chunk) == n:
                yield tuple(chunk)
            return
        yield tuple(chunk)


def _flatten_cubic(x0, y0, x1, y1, x2, y2, x3, y3, tol, out, depth=0):
    """Adaptive de Casteljau subdivision for a cubic Bezier."""
    # Flatness test: distance from control points to the chord.
    dx = x3 - x0
    dy = y3 - y0
    d1 = abs((x1 - x0) * dy - (y1 - y0) * dx)
    d2 = abs((x2 - x0) * dy - (y2 - y0) * dx)
    seg_len = math.hypot(dx, dy) or 1.0
    flatness = (d1 + d2) / seg_len
    if (flatness <= tol and depth >= 1) or depth > 18:
        out.append((x3, y3))
        return
    # Subdivide at t=0.5
    mx0, my0 = (x0 + x1) / 2, (y0 + y1) / 2
    mx1, my1 = (x1 + x2) / 2, (y1 + y2) / 2
    mx2, my2 = (x2 + x3) / 2, (y2 + y3) / 2
    mx3, my3 = (mx0 + mx1) / 2, (my0 + my1) / 2
    mx4, my4 = (mx1 + mx2) / 2, (my1 + my2) / 2
    mx, my = (mx3 + mx4) / 2, (my3 + my4) / 2
    _flatten_cubic(x0, y0, mx0, my0, mx3, my3, mx, my, tol, out, depth + 1)
    _flatten_cubic(mx, my, mx4, my4, mx2, my2, x3, y3, tol, out, depth + 1)


def _flatten_quad(x0, y0, x1, y1, x2, y2, tol, out, depth=0):
    """Adaptive subdivision for a quadratic Bezier."""
    dx = x2 - x0
    dy = y2 - y0
    d = abs((x1 - x0) * dy - (y1 - y0) * dx)
    seg_len = math.hypot(dx, dy) or 1.0
    if (d / seg_len <= tol and depth >= 1) or depth > 18:
        out.append((x2, y2))
        return
    mx0, my0 = (x0 + x1) / 2, (y0 + y1) / 2
    mx1, my1 = (x1 + x2) / 2, (y1 + y2) / 2
    mx, my = (mx0 + mx1) / 2, (my0 + my1) / 2
    _flatten_quad(x0, y0, mx0, my0, mx, my, tol, out, depth + 1)
    _flatten_quad(mx, my, mx1, my1, x2, y2, tol, out, depth + 1)


def _arc_to_cubics(x0, y0, rx_, ry_, x_rot, large, sweep, x1, y1, tol, out):
    """Flatten an SVG arc by converting to cubic Bezier segments."""
    # If endpoints are identical, arc is omitted.
    if x0 == x1 and y0 == y1:
        return
    rx_ = abs(rx_)
    ry_ = abs(ry_)
    if rx_ == 0 or ry_ == 0:
        out.append((x1, y1))
        return
    phi = math.radians(x_rot)
    cos_phi, sin_phi = math.cos(phi), math.sin(phi)
    # Step 1: compute (x1', y1')
    dx = (x0 - x1) / 2
    dy = (y0 - y1) / 2
    x1p = cos_phi * dx + sin_phi * dy
    y1p = -sin_phi * dx + cos_phi * dy
    # Correct radii
    r2 = (x1p * x1p) / (rx_ * rx_) + (y1p * y1p) / (ry_ * ry_)
    if r2 > 1:
        f = math.sqrt(r2)
        rx_ *= f
        ry_ *= f
    # Step 2: compute (cx', cy')
    sign = -1 if large == sweep else 1
    num = rx_ * rx_ * ry_ * ry_ - rx_ * rx_ * y1p * y1p - ry_ * ry_ * x1p * x1p
    den = rx_ * rx_ * y1p * y1p + ry_ * ry_ * x1p * x1p
    coef = math.sqrt(max(0, num / den)) if den else 0
    cxp = sign * coef * (rx_ * y1p / ry_)
    cyp = sign * coef * (-ry_ * x1p / rx_)
    # Step 3: compute (cx, cy)
    cx_ = cos_phi * cxp - sin_phi * cyp + (x0 + x1) / 2
    cy_ = sin_phi * cxp + cos_phi * cyp + (y0 + y1) / 2
    # Step 4: compute theta1 and delta_theta
    def angle(ux, uy, vx, vy):
        dot = ux * vx + uy * vy
        len_u = math.hypot(ux, uy)
        len_v = math.hypot(vx, vy)
        c = max(-1, min(1, dot / (len_u * len_v)))
        a = math.acos(c)
        if ux * vy - uy * vx < 0:
            a = -a
        return a
    theta1 = angle(1, 0, (x1p - cxp) / rx_, (y1p - cyp) / ry_)
    delta = angle((x1p - cxp) / rx_, (y1p - cyp) / ry_,
                  (-x1p - cxp) / rx_, (-y1p - cyp) / ry_)
    if not sweep and delta > 0:
        delta -= 2 * math.pi
    elif sweep and delta < 0:
        delta += 2 * math.pi
    # Split into segments of <= 90°
    n_segs = max(1, int(math.ceil(abs(delta) / (math.pi / 2))))
    seg_delta = delta / n_segs
    t = math.tan(seg_delta / 2)
    alpha = math.sin(seg_delta) * (math.sqrt(4 + 3 * t * t) - 1) / 3
    cos_t1, sin_t1 = math.cos(theta1), math.sin(theta1)
    for i in range(n_segs):
        theta1_i = theta1 + i * seg_delta
        cos1, sin1 = math.cos(theta1_i), math.sin(theta1_i)
        cos2, sin2 = math.cos(theta1_i + seg_delta), math.sin(theta1_i + seg_delta)
        # control point 1
        px1 = cx_ + rx_ * (cos1 - alpha * sin1) * cos_phi \
            - ry_ * (sin1 + alpha * cos1) * sin_phi
        py1 = cy_ + rx_ * (cos1 - alpha * sin1) * sin_phi \
            + ry_ * (sin1 + alpha * cos1) * cos_phi
        # control point 2
        px2 = cx_ + rx_ * (cos2 + alpha * sin2) * cos_phi \
            - ry_ * (sin2 - alpha * cos2) * sin_phi
        py2 = cy_ + rx_ * (cos2 + alpha * sin2) * sin_phi \
            + ry_ * (sin2 - alpha * cos2) * cos_phi
        # end point
        ex = cx_ + rx_ * cos2 * cos_phi - ry_ * sin2 * sin_phi
        ey = cy_ + rx_ * cos2 * sin_phi + ry_ * sin2 * cos_phi
        _flatten_cubic(x0 if i == 0 else out[-1][0],
                       y0 if i == 0 else out[-1][1],
                       px1, py1, px2, py2, ex, ey, tol, out)


# ---------------------------------------------------------------------------
# Marker (arrowhead) handling
# ---------------------------------------------------------------------------
@dataclass
class MarkerDef:
    """Parsed <marker> definition."""
    id: str
    width: float = 10
    height: float = 8
    ref_x: float = 9
    ref_y: float = 4
    fill: str = "#000000"
    points: list = field(default_factory=list)  # polygon points


def _parse_markers(root: ET.Element) -> dict[str, MarkerDef]:
    markers = {}
    ns = {"svg": "http://www.w3.org/2000/svg"}
    for m in root.iter():
        tag = _local(m.tag)
        if tag == "marker":
            mid = m.get("id", "")
            md = MarkerDef(id=mid)
            md.width = _float(m.get("markerWidth", "10"))
            md.height = _float(m.get("markerHeight", "8"))
            md.ref_x = _float(m.get("refX", "0"))
            md.ref_y = _float(m.get("refY", "0"))
            # Find the polygon inside
            for child in m:
                ct = _local(child.tag)
                if ct == "polygon":
                    pts_str = child.get("points", "")
                    md.points = _parse_points(pts_str)
                    md.fill = _normalize_color(child.get("fill", "#000000"))
                elif ct == "path":
                    # Approximate path marker as polygon
                    sub = _flatten_path(child.get("d", ""))
                    if sub:
                        md.points = sub[0][0]
                    md.fill = _normalize_color(child.get("fill", "#000000"))
            markers[mid] = md
    return markers


def _render_marker(shape_collection, marker: MarkerDef, x, y, ux, uy, scale, offset_x, offset_y):
    """Render an arrowhead marker as a freeform triangle at endpoint (x,y).

    The marker is oriented so its x-axis aligns with direction (ux, uy).
    """
    if not marker.points:
        return
    # Perpendicular to direction
    nx, ny = -uy, ux
    # Transform each marker-space point to world space:
    # world = endpoint + (mx - refX) * u + (my - refY) * n
    world_pts = []
    for mx, my in marker.points:
        wx = x + (mx - marker.ref_x) * ux + (my - marker.ref_y) * nx
        wy = y + (mx - marker.ref_x) * uy + (my - marker.ref_y) * ny
        world_pts.append((wx, wy))
    _add_freeform(shape_collection, world_pts, True, SvgStyle(fill=marker.fill),
                  scale, offset_x, offset_y)


# ---------------------------------------------------------------------------
# Shape creation
# ---------------------------------------------------------------------------
def _local(tag: str) -> str:
    return tag.split("}")[-1].lower() if "}" in tag else tag.lower()


def _parse_points(s: str) -> list[tuple[float, float]]:
    vals = re.split(r"[\s,]+", s.strip())
    vals = [v for v in vals if v]
    pts = []
    for i in range(0, len(vals) - 1, 2):
        try:
            pts.append((float(vals[i]), float(vals[i + 1])))
        except ValueError:
            continue
    return pts


def _apply_fill(shape, style: SvgStyle, config: PptxConfig):
    """Apply fill from SvgStyle to a pptx shape."""
    fill = shape.fill
    fill_color = _normalize_color(style.fill, config.default_fill)
    eff_opacity = _effective_opacity(style)
    if fill_color == "none":
        fill.background()
    else:
        fill.solid()
        rgb = _hex_to_rgbcolor(fill_color)
        if rgb:
            fill.fore_color.rgb = rgb
        if eff_opacity < 1.0:
            _set_fill_alpha(shape, eff_opacity)


def _set_fill_alpha(shape, alpha: float):
    """Set fill transparency via direct XML manipulation (python-pptx lacks API)."""
    alpha_pct = int(round(alpha * 100000))
    sp = shape.fill._xPr  # spPr or ln element
    srgb = sp.find(qn("a:solidFill") + "/" + qn("a:srgbClr"))
    if srgb is None:
        solid = sp.find(qn("a:solidFill"))
        if solid is not None:
            srgb = solid.find(qn("a:srgbClr"))
    if srgb is not None:
        alpha_el = srgb.find(qn("a:alpha"))
        if alpha_el is None:
            alpha_el = srgb.makeelement(qn("a:alpha"), {})
            srgb.append(alpha_el)
        alpha_el.set("val", str(alpha_pct))


def _apply_line(shape, style: SvgStyle, config: PptxConfig):
    """Apply stroke from SvgStyle to a pptx shape's line."""
    line = shape.line
    stroke = _normalize_color(style.stroke, config.default_stroke)
    if stroke == "none":
        line.fill.background()
        return
    rgb = _hex_to_rgbcolor(stroke)
    if rgb:
        line.color.rgb = rgb
    line.width = Emu(px(style.stroke_width))
    if style.stroke_dasharray:
        _set_line_dash(shape, style.stroke_dasharray, style.stroke_width)


def _set_line_dash(shape, dasharray: str, width: float):
    """Set a custom dash pattern via XML (python-pptx has limited dash support)."""
    # Parse "6,3" or "4 3" → dash, gap lengths in px
    parts = [float(x) for x in re.split(r"[\s,]+", dasharray.strip()) if x]
    if not parts:
        return
    # Build prstDash val from common patterns, else custom
    ln = shape.line._get_or_add_ln()
    # Remove existing dash
    for old in ln.findall(qn("a:prstDash")):
        ln.remove(old)
    for old in ln.findall(qn("a:custDash")):
        ln.remove(old)
    if len(parts) == 2 and abs(parts[0] / max(parts[1], 0.01) - 2.0) < 0.3:
        prst = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(prst)
    else:
        # Custom dash
        cust = ln.makeelement(qn("a:custDash"), {})
        for i, p in enumerate(parts):
            d_len = int(p * EMU_PER_PX)
            tag = qn("a:ds") if i % 2 == 0 else qn("a:g")
            el = cust.makeelement(tag, {"d": str(max(d_len, 1)), "sp": str(max(d_len, 1))})
            cust.append(el)
        ln.append(cust)


def _disable_shadow(shape):
    """Disable inherited shadow (python-pptx autoShape default has shadow)."""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _estimate_text_width(text: str, font_size: float, weight: str) -> float:
    """Estimate text width in px (matches SVGDrawer / evaluator metric)."""
    coef = 0.62 if weight == "bold" else 0.55
    return sum(font_size * (1.0 if ord(c) > 0x2E80 else coef) for c in text)


def _create_text(shapes, el: ET.Element, transform, style: SvgStyle,
                 scale, offset_x, offset_y):
    """Create a PowerPoint text box from an SVG <text> element."""
    x = _float(el.get("x", "0"))
    y = _float(el.get("y", "0"))
    tx, ty = transform_point(transform, (x, y))
    # Concatenate direct text + <tspan> children (itertext handles both).
    content = "".join(el.itertext())
    content = html.unescape(content)
    if not content.strip():
        return

    fs = style.font_size
    w_est = _estimate_text_width(content, fs, style.font_weight)
    h_est = fs * 1.3

    anchor = style.text_anchor
    baseline = style.dominant_baseline

    if anchor == "middle":
        left = tx - w_est / 2
        align = PP_ALIGN.CENTER
    elif anchor == "end":
        left = tx - w_est
        align = PP_ALIGN.RIGHT
    else:
        left = tx
        align = PP_ALIGN.LEFT

    if baseline in ("central", "middle"):
        top = ty - h_est / 2
        vert = MSO_ANCHOR.MIDDLE
    elif baseline in ("hanging", "text-before-edge"):
        top = ty
        vert = MSO_ANCHOR.TOP
    else:
        top = ty - fs
        vert = MSO_ANCHOR.BOTTOM

    # Add padding to avoid clipping
    pad = fs * 0.3
    left -= pad
    top -= pad * 0.3
    w_est += pad * 2
    h_est += pad * 0.6

    box_left = offset_x + px(left * scale)
    box_top = offset_y + px(top * scale)
    box_w = px(w_est * scale)
    box_h = px(h_est * scale)

    tx_shape = shapes.add_textbox(box_left, box_top, box_w, box_h)
    tx_shape.text_frame.word_wrap = False
    tx_shape.text_frame.margin_left = 0
    tx_shape.text_frame.margin_right = 0
    tx_shape.text_frame.margin_top = 0
    tx_shape.text_frame.margin_bottom = 0
    tx_shape.text_frame.vertical_anchor = vert

    p = tx_shape.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(fs * scale * (72 / 96))  # px→pt: 1pt = 1.333px at 96dpi
    run.font.bold = (style.font_weight == "bold")
    run.font.italic = (style.font_style == "italic")
    fill_color = _normalize_color(style.fill, "#000000")
    rgb = _hex_to_rgbcolor(fill_color)
    if rgb:
        run.font.color.rgb = rgb

    # Disable textbox border/fill
    _disable_shadow(tx_shape)
    tx_shape.fill.background()
    tx_shape.line.fill.background()


def _add_freeform(shapes, points, closed, style, scale, offset_x, offset_y):
    """Add a freeform polygon/polyline to the shapes collection."""
    if len(points) < 2:
        return None
    # Transform points to absolute EMU
    emu_pts = [(offset_x + px(p[0] * scale), offset_y + px(p[1] * scale))
               for p in points]
    try:
        builder = shapes.build_freeform(emu_pts[0][0], emu_pts[0][1])
        builder.add_line_segments(emu_pts[1:], close=closed)
        shape = builder.convert_to_shape()
    except Exception:
        # Fallback: if freeform fails, skip this shape
        return None

    _apply_fill(shape, style, _noop_config())
    _apply_line(shape, style, _noop_config())
    _disable_shadow(shape)
    return shape


_noop_cfg = None


def _noop_config():
    global _noop_cfg
    if _noop_cfg is None:
        _noop_cfg = PptxConfig()
    return _noop_cfg


# ---------------------------------------------------------------------------
# Main converter
# ---------------------------------------------------------------------------
class _Converter:
    """Walks the SVG element tree and creates PPTX shapes."""

    def __init__(self, config: PptxConfig):
        self.config = config
        self.markers: dict[str, MarkerDef] = {}

    def convert(self, svg_content: str) -> Presentation:
        """Parse SVG string, return a Presentation with one slide."""
        # Strip XML namespace prefixes for easier parsing
        svg_clean = svg_content
        root = ET.fromstring(svg_content)
        # Register namespace
        ET.register_namespace("", "http://www.w3.org/2000/svg")

        # Get SVG dimensions
        vb = root.get("viewBox")
        if vb:
            parts = [float(x) for x in re.split(r"[\s,]+", vb.strip())]
            vb_x, vb_y, vb_w, vb_h = parts[0], parts[1], parts[2], parts[3]
        else:
            vb_x = vb_y = 0
            vb_w = _float(root.get("width", "1200"))
            vb_h = _float(root.get("height", "800"))

        # Parse markers from defs
        self.markers = _parse_markers(root)

        # Create presentation
        prs = Presentation()
        prs.slide_width = Emu(int(self.config.slide_w * 914400))
        prs.slide_height = Emu(int(self.config.slide_h * 914400))
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if self.config.mode == "image":
            self._add_image(svg_content, slide, vb_w, vb_h, prs)
        else:
            # Calculate fit-to-slide scale
            margin_px = self.config.margin * 96
            avail_w = self.config.slide_w * 96 - 2 * margin_px
            avail_h = self.config.slide_h * 96 - 2 * margin_px
            fit_scale = min(avail_w / vb_w, avail_h / vb_h)
            scale = fit_scale * self.config.scale

            # Center the SVG content
            rendered_w = vb_w * scale
            rendered_h = vb_h * scale
            offset_x = px((self.config.slide_w * 96 - rendered_w) / 2 - vb_x * scale)
            offset_y = px((self.config.slide_h * 96 - rendered_h) / 2 - vb_y * scale)

            # Walk elements
            for child in root:
                self._walk(child, IDENTITY, None, slide.shapes, scale,
                           offset_x, offset_y)

        return prs

    def _walk(self, el, parent_transform, parent_style, shapes, scale, ox, oy):
        tag = _local(el.tag)
        if tag in ("defs", "marker", "clippath", "filter",
                    "lineargradient", "radialgradient", "pattern", "symbol"):
            return  # skip non-rendering elements

        own_t = parse_transform(el.get("transform", ""))
        transform = multiply_matrix(parent_transform, own_t)
        style = _parse_style(el, parent_style)

        if tag == "g":
            for child in el:
                self._walk(child, transform, style, shapes, scale, ox, oy)

        elif tag == "rect":
            self._make_rect(el, transform, style, shapes, scale, ox, oy)

        elif tag in ("circle", "ellipse"):
            self._make_oval(el, tag, transform, style, shapes, scale, ox, oy)

        elif tag == "line":
            self._make_line(el, transform, style, shapes, scale, ox, oy)

        elif tag in ("polygon", "polyline"):
            self._make_polygon(el, tag, transform, style, shapes, scale, ox, oy)

        elif tag == "path":
            self._make_path(el, transform, style, shapes, scale, ox, oy)

        elif tag == "text":
            _create_text(shapes, el, transform, style, scale, ox, oy)

        elif tag == "use":
            # Resolve <use href="#id"> — find the referenced element and render it
            href = el.get("href") or el.get("{http://www.w3.org/1999/xlink}href")
            if href:
                ref_id = href.lstrip("#")
                # Would need a symbol/defs registry; skip for now
                pass

    def _make_rect(self, el, transform, style, shapes, scale, ox, oy):
        x = _float(el.get("x", "0"))
        y = _float(el.get("y", "0"))
        w = _float(el.get("width", "0"))
        h = _float(el.get("height", "0"))
        rx_ = _float(el.get("rx", "0"))
        ry_ = _float(el.get("ry", str(rx_)))
        if w <= 0 or h <= 0:
            return
        # Apply transform to all four corners, get AABB
        corners = [(x, y), (x + w, y), (x, y + h), (x + w, y + h)]
        pts = [transform_point(transform, c) for c in corners]
        xs = [p[0] for p in pts]
        ys = [p[1] for p in pts]
        ax, ay = min(xs), min(ys)
        aw, ah = max(xs) - ax, max(ys) - ay

        # Check if transform is axis-aligned (no rotation/skew)
        is_axis_aligned = (abs(transform[1]) < 1e-6 and abs(transform[2]) < 1e-6)

        if rx_ > 0 or ry_ > 0:
            shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Emu(ox + px(ax * scale)),
                                     Emu(oy + px(ay * scale)),
                                     Emu(px(aw * scale)),
                                     Emu(px(ah * scale)))
            # Adjust corner radius (adjustment value 0-0.5)
            try:
                adj = min(rx_, ry_) / min(w, h)
                shape.adjustments[0] = min(0.5, adj)
            except Exception:
                pass
        else:
            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Emu(ox + px(ax * scale)),
                                     Emu(oy + px(ay * scale)),
                                     Emu(px(aw * scale)),
                                     Emu(px(ah * scale)))

        # Apply rotation uniformly to both rect variants.
        # For a rectangle, the AABB center == rotated-rect center (by symmetry),
        # and PPTX rotates around the shape center, so this is geometrically exact
        # for rotation+scale+translate transforms.
        if not is_axis_aligned:
            angle = math.degrees(math.atan2(transform[1], transform[0]))
            shape.rotation = angle

        _apply_fill(shape, style, self.config)
        _apply_line(shape, style, self.config)
        _disable_shadow(shape)

    def _make_oval(self, el, tag, transform, style, shapes, scale, ox, oy):
        cx_ = _float(el.get("cx", "0"))
        cy_ = _float(el.get("cy", "0"))
        if tag == "circle":
            r = _float(el.get("r", "0"))
            rx_ = ry_ = r
        else:
            rx_ = _float(el.get("rx", "0"))
            ry_ = _float(el.get("ry", "0"))
        if rx_ <= 0 or ry_ <= 0:
            return
        x = cx_ - rx_
        y = cy_ - ry_
        w = 2 * rx_
        h = 2 * ry_
        corners = [(x, y), (x + w, y), (x, y + h), (x + w, y + h)]
        pts = [transform_point(transform, c) for c in corners]
        xs = [p[0] for p in pts]
        ys = [p[1] for p in pts]
        ax, ay = min(xs), min(ys)
        aw, ah = max(xs) - ax, max(ys) - ay

        shape = shapes.add_shape(MSO_SHAPE.OVAL,
                                 Emu(ox + px(ax * scale)),
                                 Emu(oy + px(ay * scale)),
                                 Emu(px(aw * scale)),
                                 Emu(px(ah * scale)))
        # Apply rotation if transform is non-axis-aligned (same reasoning as rect).
        if abs(transform[1]) > 1e-6 or abs(transform[2]) > 1e-6:
            angle = math.degrees(math.atan2(transform[1], transform[0]))
            shape.rotation = angle
        _apply_fill(shape, style, self.config)
        _apply_line(shape, style, self.config)
        _disable_shadow(shape)

    def _make_line(self, el, transform, style, shapes, scale, ox, oy):
        x1 = _float(el.get("x1", "0"))
        y1 = _float(el.get("y1", "0"))
        x2 = _float(el.get("x2", "0"))
        y2 = _float(el.get("y2", "0"))
        p1 = transform_point(transform, (x1, y1))
        p2 = transform_point(transform, (x2, y2))

        conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Emu(ox + px(p1[0] * scale)),
                                    Emu(oy + px(p1[1] * scale)),
                                    Emu(ox + px(p2[0] * scale)),
                                    Emu(oy + px(p2[1] * scale)))
        _disable_shadow(conn)
        stroke = _normalize_color(style.stroke, self.config.default_stroke)
        if stroke != "none":
            rgb = _hex_to_rgbcolor(stroke)
            if rgb:
                conn.line.color.rgb = rgb
        conn.line.width = Emu(px(style.stroke_width))
        if style.stroke_dasharray:
            _set_line_dash(conn, style.stroke_dasharray, style.stroke_width)

        # Render marker-end (arrowhead) if present
        marker_ref = el.get("marker-end", "")
        if marker_ref:
            mid = re.search(r"url\(#([^)]+)\)", marker_ref)
            if mid and mid.group(1) in self.markers:
                marker = self.markers[mid.group(1)]
                dx, dy = p2[0] - p1[0], p2[1] - p1[1]
                seglen = math.hypot(dx, dy) or 1.0
                ux, uy = dx / seglen, dy / seglen
                _render_marker(shapes, marker, p2[0], p2[1], ux, uy, scale, ox, oy)

        marker_start = el.get("marker-start", "")
        if marker_start:
            mid = re.search(r"url\(#([^)]+)\)", marker_start)
            if mid and mid.group(1) in self.markers:
                marker = self.markers[mid.group(1)]
                dx, dy = p1[0] - p2[0], p1[1] - p2[1]
                seglen = math.hypot(dx, dy) or 1.0
                ux, uy = dx / seglen, dy / seglen
                _render_marker(shapes, marker, p1[0], p1[1], ux, uy, scale, ox, oy)

    def _make_polygon(self, el, tag, transform, style, shapes, scale, ox, oy):
        pts_str = el.get("points", "")
        local_pts = _parse_points(pts_str)
        if len(local_pts) < 2:
            return
        world_pts = [transform_point(transform, p) for p in local_pts]
        closed = (tag == "polygon")
        _add_freeform(shapes, world_pts, closed, style, scale, ox, oy)

    def _make_path(self, el, transform, style, shapes, scale, ox, oy):
        d = el.get("d", "")
        if not d.strip():
            return
        subpaths = _flatten_path(d, self.config.curve_tolerance)
        for pts, closed in subpaths:
            if len(pts) < 2:
                continue
            world_pts = [transform_point(transform, p) for p in pts]
            _add_freeform(shapes, world_pts, closed, style, scale, ox, oy)

        # Render marker-end for paths too
        marker_ref = el.get("marker-end", "")
        if marker_ref and subpaths:
            mid = re.search(r"url\(#([^)]+)\)", marker_ref)
            if mid and mid.group(1) in self.markers:
                marker = self.markers[mid.group(1)]
                last_pts = subpaths[-1][0]
                if len(last_pts) >= 2:
                    p1 = transform_point(transform, last_pts[-2])
                    p2 = transform_point(transform, last_pts[-1])
                    dx, dy = p2[0] - p1[0], p2[1] - p1[1]
                    seglen = math.hypot(dx, dy) or 1.0
                    ux, uy = dx / seglen, dy / seglen
                    _render_marker(shapes, marker, p2[0], p2[1], ux, uy, scale, ox, oy)

    def _add_image(self, svg_content, slide, vb_w, vb_h, prs):
        """Rasterize SVG to PNG via rsvg-convert and embed as a picture."""
        png_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".svg", delete=False, mode="w") as sf:
                sf.write(svg_content)
                svg_file = sf.name
            png_path = svg_file.replace(".svg", ".png")
            dpi = self.config.image_dpi
            scale = dpi / 96
            subprocess.run(
                ["rsvg-convert", "-d", str(dpi), "-p", str(dpi),
                 svg_file, "-o", png_path],
                check=True, capture_output=True)
            # Compute placement
            slide_w = self.config.slide_w * 914400
            slide_h = self.config.slide_h * 914400
            margin = self.config.margin * 914400
            avail_w = slide_w - 2 * margin
            avail_h = slide_h - 2 * margin
            # Aspect ratio from SVG
            ratio = vb_w / vb_h
            target_w = avail_w
            target_h = int(target_w / ratio)
            if target_h > avail_h:
                target_h = avail_h
                target_w = int(target_h * ratio)
            left = int((slide_w - target_w) / 2)
            top = int((slide_h - target_h) / 2)
            slide.shapes.add_picture(png_path, Emu(left), Emu(top),
                                     Emu(int(target_w)), Emu(int(target_h)))
        finally:
            for f in ([png_path] if png_path else []):
                try:
                    os.unlink(f)
                except OSError:
                    pass
            try:
                os.unlink(svg_file)
            except (OSError, NameError):
                pass


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------
def svg_to_pptx(
    svg: Union[str, Path],
    pptx_path: Union[str, Path],
    config: Optional[PptxConfig] = None,
) -> Presentation:
    """Convert SVG content to a PowerPoint file.

    Args:
        svg:        SVG content as a string, or a path to an .svg file.
        pptx_path:  Output .pptx file path.
        config:     Optional PptxConfig. Defaults to 16:9, shapes mode.

    Returns:
        The pptx Presentation object (also saved to pptx_path).
    """
    cfg = config or PptxConfig()
    # Determine if svg is a file path or SVG string
    svg_str = svg
    if isinstance(svg, (str, Path)):
        sp = str(svg)
        if os.path.isfile(sp) and sp.lower().endswith(".svg"):
            with open(sp, "r", encoding="utf-8") as f:
                svg_str = f.read()
    # If svg_to_pptx is called with a pure SVG string that happens to not exist
    # as a file, use it directly (it's already the SVG content).
    converter = _Converter(cfg)
    prs = converter.convert(svg_str)
    out = Path(pptx_path).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return prs


def add_svg_to_slide(
    svg: Union[str, Path],
    slide,
    x: float = 0,
    y: float = 0,
    scale: float = 1.0,
    config: Optional[PptxConfig] = None,
):
    """Add SVG shapes to an existing slide (in-place).

    Args:
        svg:    SVG content string or .svg file path.
        slide:  pptx Slide object to add shapes to.
        x, y:   Top-left placement in inches.
        scale:  Scale factor (1.0 = original pixel size mapped to EMU).
        config: Optional PptxConfig (mode/curve_tolerance used).
    """
    cfg = config or PptxConfig()
    svg_str = svg
    if isinstance(svg, (str, Path)):
        sp = str(svg)
        if os.path.isfile(sp) and sp.lower().endswith(".svg"):
            with open(sp, "r", encoding="utf-8") as f:
                svg_str = f.read()
    root = ET.fromstring(svg_str)
    ET.register_namespace("", "http://www.w3.org/2000/svg")
    converter = _Converter(cfg)
    converter.markers = _parse_markers(root)
    vb = root.get("viewBox")
    if vb:
        parts = [float(x2) for x2 in re.split(r"[\s,]+", vb.strip())]
        vb_x, vb_y = parts[0], parts[1]
    else:
        vb_x = vb_y = 0
    offset_x = px(x * 96 - vb_x * scale)
    offset_y = px(y * 96 - vb_y * scale)
    for child in root:
        converter._walk(child, IDENTITY, None, slide.shapes, scale,
                        offset_x, offset_y)


def save_pptx(drawer, pptx_path, config: Optional[PptxConfig] = None):
    """Convert an SVGDrawer's rendered output to PPTX.

    Convenience wrapper: drawer.render() → svg_to_pptx().
    """
    return svg_to_pptx(drawer.render(), pptx_path, config)
